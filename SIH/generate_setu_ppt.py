"""
SETU - SIH 2026 Idea Presentation
Uses the EXACT template PPTX and replaces content with SETU project data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ---- PATHS ----
TEMPLATE = r"d:\Testing\Research\SIH\Sih-ppt-template-2025-pdf-download.pptx"
OUT_DIR  = r"d:\Testing\Research\SIH"
OUT_PPTX = os.path.join(OUT_DIR, "SETU_SIH2026_Idea_Presentation.pptx")

# ---- COLORS ----
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
LIGHT_CYAN = RGBColor(0xE0, 0xF7, 0xFA)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
PINK = RGBColor(0xE9, 0x1E, 0x63)
LIGHT_PINK = RGBColor(0xFC, 0xE4, 0xEC)
BLUE = RGBColor(0x21, 0x96, 0xF3)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
DEEP_ORANGE = RGBColor(0xFF, 0x57, 0x22)
PURPLE = RGBColor(0x9C, 0x27, 0xB0)
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)
TEAL = RGBColor(0x00, 0x96, 0x88)
RED = RGBColor(0xF4, 0x43, 0x36)
SAFFRON = RGBColor(0xFF, 0x99, 0x33)
IND_GREEN = RGBColor(0x13, 0x88, 0x08)
NAVY = RGBColor(0x00, 0x0E, 0x4A)
CORAL = RGBColor(0xFF, 0x6B, 0x6B)
AMBER = RGBColor(0xFF, 0xC1, 0x07)
OLIVE_BG = RGBColor(0xE8, 0xE0, 0xD0)
OLIVE_DARK = RGBColor(0x6B, 0x6B, 0x4E)

# ---- HELPERS ----

def clear_tf(tf):
    """Clear all content from a text frame, keeping one empty paragraph."""
    # Remove all paragraphs except the first
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    body = tf._txBody
    paras = body.findall('.//a:p', nsmap)
    for p in paras[1:]:
        body.remove(p)
    # Clear the first paragraph
    first_p = paras[0]
    for child in list(first_p):
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
        if tag != 'pPr':  # keep paragraph properties
            first_p.remove(child)

def rewrite_tf(tf, lines, size=14, bold=False, color=BLACK, align=None):
    """Rewrite a text frame with multiple lines of text."""
    clear_tf(tf)
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align:
            p.alignment = align
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        
        if isinstance(line_text, list):
            # Rich text: list of (text, bold, color) tuples
            for txt, b, c in line_text:
                run = p.add_run()
                run.text = txt
                run.font.size = Pt(size)
                run.font.bold = b
                run.font.color.rgb = c
                run.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = line_text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"

def textbox(slide, left, top, w, h, text="", size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def rounded_rect(slide, left, top, w, h, fill_color=None, border_color=None, bw=Pt(1.5)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = bw
        shp.line.fill.solid()
    return shp

def rect(slide, left, top, w, h, fill_color=None, border_color=None, bw=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = bw
        shp.line.fill.solid()
    return shp

def down_arrow(slide, left, top, w, h, fill_color=DARK_GRAY):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    return shp

def right_arrow_shape(slide, left, top, w, h, fill_color=DARK_GRAY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    return shp

def diamond(slide, left, top, w, h, fill_color=None, border_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(2)
        shp.line.fill.solid()
    return shp

def chevron(slide, left, top, w, h, fill_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    return shp

def oval(slide, left, top, w, h, fill_color=None, border_color=None, bw=Pt(2)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = bw
        shp.line.fill.solid()
    return shp

def set_shape_text(shape, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"

def flow_box(slide, left, top, w, h, text, fill_color, text_color=WHITE, size=10):
    shp = rounded_rect(slide, left, top, w, h, fill_color=fill_color)
    set_shape_text(shp, text, size=size, bold=True, color=text_color)
    return shp

def add_card(slide, left, top, w, h, title, items, fill_color, title_color=DARK_GRAY, border_color=None, tsz=15, isz=11):
    card = rounded_rect(slide, left, top, w, h, fill_color=fill_color, border_color=border_color)
    textbox(slide, left + 0.12, top + 0.08, w - 0.24, 0.35, title, size=tsz, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.45), Inches(w - 0.24), Inches(h - 0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = "  " + item
        run.font.size = Pt(isz)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"


# ====================================================================
#  LOAD TEMPLATE & MODIFY
# ====================================================================

print("Loading template...")
prs = Presentation(TEMPLATE)

def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

# ============ SLIDE 1: TITLE ============
print("Slide 1 - Title...")
s1 = prs.slides[0]

# Update "2025" -> "2026" in the title
sh = find_shape(s1, 'Title 7')
if sh:
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            r.text = r.text.replace('2025', '2026')

# Update subtitle
sh = find_shape(s1, 'Subtitle 3')
if sh:
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if 'TITLE PAGE' in r.text:
                r.text = r.text.replace('TITLE PAGE', 'Crowdsourced Road Defect Intelligence for India')

# Update the details textbox
sh = find_shape(s1, 'TextBox 9')
if sh:
    tf = sh.text_frame
    details = [
        [("  Problem Statement ID - ", False, BLACK), ("SIH 2026 (TBD)", True, BLACK)],
        [("  Problem Statement Title - ", False, BLACK), ("Crowdsourced Road Defect Intelligence", True, BLACK)],
        [("  Theme - ", False, BLACK), ("Smart Automation", True, BLACK)],
        [("  PS Category - ", False, BLACK), ("Software", True, BLACK)],
        [("  Team ID - ", False, BLACK), ("TBD", True, BLACK)],
        [("  Team Name - ", False, BLACK), ("Team SETU", True, BLACK)],
    ]
    rewrite_tf(tf, details, size=16)


# ============ SLIDE 2: PROPOSED SOLUTION ============
print("Slide 2 - Proposed Solution...")
s2 = prs.slides[1]

# Update team oval
for sh in s2.shapes:
    if 'Oval' in sh.name and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('Your Team Name', 'Team SETU')

# Update title
sh = find_shape(s2, 'Title 1')
if sh:
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if 'IDEA TITLE' in r.text:
                r.text = r.text.replace('IDEA TITLE', 'Crowdsourced Road Defect Intelligence')

# Clear the TextBox 8 placeholder
sh = find_shape(s2, 'TextBox 8')
if sh:
    clear_tf(sh.text_frame)

# Update footer
sh = find_shape(s2, 'Footer Placeholder 6')
if sh:
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            r.text = '@SIH Idea submission - SETU'

# Add 3-column cards
add_card(s2, 0.15, 1.35, 3.35, 4.2,
         "Problem:", [
             "9,438 pothole deaths (2020-24), rising 53% in 5 years",
             "6+ deaths every day from a hole in the road",
             "States report ZERO crashes -- data is broken, not just roads",
             "Manual inspection takes weeks for 1000+ km",
             "Citizen apps are reactive & gameable",
         ], LIGHT_CYAN, title_color=CYAN, border_color=CYAN, tsz=14, isz=10)

add_card(s2, 3.65, 1.35, 3.35, 4.2,
         "Our Solution:", [
             "SETU: Sensor-Enabled Tracking of Urban-road-damage",
             "Turn 1.2Cr gig workers' phones into a free road inspection network",
             "3-Layer Cost-Gated Escalation: Sensors > Clustering > Vision AI",
             "Cost ~Rs.3 per confirmed pothole vs Rs.17,693 manual audit",
         ], LIGHT_GREEN, title_color=DARK_GREEN, border_color=GREEN, tsz=14, isz=10)

add_card(s2, 7.15, 1.35, 3.15, 4.2,
         "Uniqueness:", [
             "Nobody does all 3 layers with cost-gated escalation",
             "Per-device self-calibration across vehicles",
             "Negative Evidence: certifies good roads too",
             "Repair verification by silence -- audit trail for municipalities",
             "DPDP-compliant from day 1",
         ], LIGHT_PINK, title_color=PINK, border_color=PINK, tsz=14, isz=10)

# Vertical Flowchart (right side)
fc_x = 10.55
fc_w = 2.45
fc_h = 0.6
flow_box(s2, fc_x, 1.35, fc_w, fc_h, "Phone Sensors\n(Accel+Gyro @100Hz)", SAFFRON, size=9)
down_arrow(s2, fc_x+fc_w/2-0.1, 1.95, 0.2, 0.25, DARK_GRAY)
flow_box(s2, fc_x, 2.25, fc_w, fc_h, "On-Device TFLite\n6-class classifier", BLUE, size=9)
down_arrow(s2, fc_x+fc_w/2-0.1, 2.85, 0.2, 0.25, DARK_GRAY)
flow_box(s2, fc_x, 3.15, fc_w, fc_h, "Server Clustering\nDBSCAN + Bayesian", TEAL, size=9)
down_arrow(s2, fc_x+fc_w/2-0.1, 3.75, 0.2, 0.25, DARK_GRAY)
flow_box(s2, fc_x, 4.05, fc_w, fc_h, "YOLO Vision\nConfirms @ 0.1% spots", DEEP_ORANGE, size=9)
down_arrow(s2, fc_x+fc_w/2-0.1, 4.65, 0.2, 0.25, DARK_GRAY)
shp_ok = rounded_rect(s2, fc_x, 4.95, 1.1, 0.5, fill_color=GREEN)
set_shape_text(shp_ok, "Confirmed", size=10, bold=True, color=WHITE)
shp_no = rounded_rect(s2, fc_x+1.3, 4.95, 1.1, 0.5, fill_color=RED)
set_shape_text(shp_no, "Rejected", size=10, bold=True, color=WHITE)


# ============ SLIDE 3: TECHNICAL APPROACH ============
print("Slide 3 - Technical Approach...")
s3 = prs.slides[2]

for sh in s3.shapes:
    if 'Oval' in sh.name and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('Your Team Name', 'Team SETU')
    if sh.name == 'TextBox 8':
        clear_tf(sh.text_frame)
    if sh.name == 'Footer Placeholder 6':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = '@SIH Idea submission - SETU'

# Left: System Architecture Flow
boxes_data = [
    ("Android App / SDK\n(Kotlin, 200KB AAR)", PINK),
    ("Sensor Capture\n(Accel+Gyro @100Hz)", CYAN),
    ("On-Device TFLite ML\n(1D-CNN, 120KB INT8)", ORANGE),
    ("FastAPI Backend\n(Clustering + Brain)", BLUE),
    ("YOLO Vision Pipeline\n(ByteTrack + Depth)", TEAL),
    ("Municipal Dashboard\n(React + deck.gl)", DARK_GREEN),
]
y = 1.35
for i, (text, bg_col) in enumerate(boxes_data):
    flow_box(s3, 0.3, y, 2.7, 0.58, text, bg_col, WHITE, size=9)
    if i < len(boxes_data) - 1:
        down_arrow(s3, 0.3+2.7/2-0.08, y+0.58, 0.16, 0.18, DARK_GRAY)
    y += 0.58 + 0.22

# Separator
rect(s3, 3.25, 1.35, 0.02, 5.15, fill_color=DARK_GRAY)

# Right: Tech grid
textbox(s3, 3.5, 1.15, 9.3, 0.4, "TECHNOLOGY POWERING THE PLATFORM",
        size=17, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

tech = [
    [("Android App", "Kotlin, Jetpack Compose, CameraX,\nSensorManager, WorkManager, TFLite", LIGHT_CYAN, CYAN),
     ("Sensor Processing", "100Hz IMU, Band-pass filter 0.5-30Hz,\nOrientation correction, Self-calibration", LIGHT_GREEN, GREEN),
     ("On-Device ML", "1D-CNN (120KB INT8 quantised),\nRandom Forest fallback,\n6-class pothole classifier", LIGHT_ORANGE, ORANGE)],
    [("Backend API", "FastAPI, Celery, Redis Streams,\nPydantic, JWT auth, OSRM\nmap-matching", LIGHT_BLUE, BLUE),
     ("Database", "PostgreSQL 16 + PostGIS 3,\nH3 hex indexing, partitioned events,\nBayesian posterior fusion", LIGHT_PURPLE, PURPLE),
     ("Vision AI", "YOLOv8/YOLO11 fine-tuned on\nRDD2022, ByteTrack dedup,\nScene classifier, Depth estimation", LIGHT_PINK, DEEP_ORANGE)],
    [("Web Dashboard", "React + TypeScript, deck.gl,\nGoogle Maps, WebSocket live\nupdates, RBAC admin portal", LIGHT_CYAN, TEAL),
     ("Infrastructure", "Docker Compose, Nginx, S3/MinIO,\nMLflow model versioning, CI/CD,\nAlembic DB migrations", LIGHT_GREEN, DARK_GREEN),
     ("Privacy & Security", "DPDP Act compliant, Argon2id,\nTOTP 2FA, No PII stored,\n7-day video retention", LIGHT_ORANGE, RED)],
]
cw, ch = 3.05, 1.3
sx, sy = 3.55, 1.65
gx, gy = 0.12, 0.1
for ri, row in enumerate(tech):
    for ci, (title, desc, bg, accent) in enumerate(row):
        x = sx + ci * (cw + gx)
        y2 = sy + ri * (ch + gy)
        rounded_rect(s3, x, y2, cw, ch, fill_color=bg, border_color=accent, bw=Pt(1.5))
        banner = rect(s3, x, y2, cw, 0.28, fill_color=accent)
        set_shape_text(banner, title, size=9, bold=True, color=WHITE)
        textbox(s3, x+0.06, y2+0.3, cw-0.12, ch-0.35, desc, size=8, color=DARK_GRAY)


# ============ SLIDE 4: FEASIBILITY AND VIABILITY ============
print("Slide 4 - Feasibility & Viability...")
s4 = prs.slides[3]

for sh in s4.shapes:
    if 'Oval' in sh.name and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('Your Team Name', 'Team SETU')
    if sh.name == 'TextBox 8':
        clear_tf(sh.text_frame)
    if sh.name == 'Footer Placeholder 6':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = '@SIH Idea submission - SETU'

# Left: Challenges
textbox(s4, 0.3, 1.15, 3.8, 0.4, "Challenges & Risks", size=16, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
challenges = [
    ("Android OEM Killers: ", "Xiaomi/Oppo kill background services aggressively"),
    ("GPS Error 27-32m: ", "Makes naive clustering impossible at highway speeds"),
    ("Phone Heterogeneity: ", "Different phones & vehicles produce wildly different signals"),
    ("Privacy & DPDP: ", "Video captures faces/plates; strict compliance mandatory"),
]
cy = 1.65
for title, desc in challenges:
    chev = chevron(s4, 0.3, cy, 3.8, 0.7, fill_color=CORAL)
    tf = chev.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r1 = tf.paragraphs[0].add_run()
    r1.text = title
    r1.font.size = Pt(9)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Calibri"
    r2 = tf.paragraphs[0].add_run()
    r2.text = desc
    r2.font.size = Pt(8)
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"
    cy += 0.85

# Center: Feasibility
feas_title = rounded_rect(s4, 4.5, 1.1, 4.3, 0.45, fill_color=BLUE)
set_shape_text(feas_title, "Feasibility", size=15, bold=True, color=WHITE)

feas = [
    ("Technical: ", "TFLite + FastAPI + PostGIS + YOLO. All open-source, proven stack.", SAFFRON),
    ("Operational: ", "SDK model = zero user acquisition. City buses as pilot fleet.", BLUE),
    ("Economic: ", "~Rs.3/pothole vs Rs.17,693 manual. Free-tier cloud for demo.", GREEN),
    ("Legal: ", "DPDP-compliant design. No PII. Consent-first video. Audit trails.", NAVY),
]
fy = 1.65
for title, desc, color in feas:
    card = rounded_rect(s4, 4.5, fy, 4.3, 0.75, border_color=color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r1 = tf.paragraphs[0].add_run()
    r1.text = title
    r1.font.size = Pt(9)
    r1.font.bold = True
    r1.font.color.rgb = color
    r1.font.name = "Calibri"
    r2 = tf.paragraphs[0].add_run()
    r2.text = desc
    r2.font.size = Pt(8)
    r2.font.color.rgb = DARK_GRAY
    r2.font.name = "Calibri"
    fy += 0.85

# Right: Mitigations
textbox(s4, 9.2, 1.15, 3.8, 0.4, "Mitigation Measures", size=16, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
mitigations = [
    "Foreground service + autostart guide; test on real Xiaomi/Oppo devices",
    "DBSCAN(eps=20m) + map-matching + inverse-variance weighting fixes GPS",
    "Per-device self-calibration in first 10 min normalises all signals",
    "No PII, hashed IDs, 7-day video retention, consent toggles, full audit trail",
]
my = 1.65
for text in mitigations:
    oval(s4, 12.4, my+0.05, 0.5, 0.5, border_color=AMBER, bw=Pt(3))
    oval(s4, 12.48, my+0.13, 0.34, 0.34, fill_color=AMBER)
    textbox(s4, 9.2, my, 3.1, 0.7, text, size=9, color=DARK_GRAY)
    my += 0.85


# ============ SLIDE 5: IMPACT AND BENEFITS ============
print("Slide 5 - Impact & Benefits...")
s5 = prs.slides[4]

for sh in s5.shapes:
    if 'Oval' in sh.name and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('Your Team Name', 'Team SETU')
    if sh.name == 'TextBox 8':
        clear_tf(sh.text_frame)
    if sh.name == 'Footer Placeholder 6':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = '@SIH Idea submission - SETU'

# Left Panel: Impact
rounded_rect(s5, 0.3, 1.35, 4.0, 5.0, fill_color=LIGHT_PINK, border_color=CORAL)
textbox(s5, 0.45, 1.4, 3.7, 0.4, "Societal Impact:", size=18, bold=True, color=DARK_GRAY)
impacts = [
    "Prevents 6+ pothole deaths per day with early detection & faster repairs",
    "Turns 1.2 crore gig workforce into free, always-on road inspection network",
    "Creates accountability: timestamped before/after data verifies contractor payments",
    "Reduces 3-5% GDP loss from road crashes (World Bank)",
    "Data-driven repair prioritisation using IRC:82-2015 standards",
]
txBox = s5.shapes.add_textbox(Inches(0.45), Inches(1.85), Inches(3.7), Inches(4.3))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(impacts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    p.space_after = Pt(5)
    run = p.add_run()
    run.text = "  " + item
    run.font.size = Pt(11)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"

# Center: Diamond arrangement
d_sz = 1.2
cx, cy2 = 5.75, 3.2
d1 = diamond(s5, cx, cy2-1.4, d_sz, d_sz, fill_color=SAFFRON)
set_shape_text(d1, "Data", size=11, bold=True, color=WHITE)
d2 = diamond(s5, cx-1.4, cy2, d_sz, d_sz, fill_color=CORAL)
set_shape_text(d2, "Detect", size=11, bold=True, color=WHITE)
dc = diamond(s5, cx, cy2, d_sz, d_sz, fill_color=TEAL)
set_shape_text(dc, "SETU", size=14, bold=True, color=WHITE)
d3 = diamond(s5, cx+1.4, cy2, d_sz, d_sz, fill_color=OLIVE_DARK)
set_shape_text(d3, "Repair", size=11, bold=True, color=WHITE)
d4 = diamond(s5, cx, cy2+1.4, d_sz, d_sz, fill_color=BLUE)
set_shape_text(d4, "Verify", size=11, bold=True, color=WHITE)

# Right Panel: Benefits
rounded_rect(s5, 9.0, 1.35, 4.0, 5.0, fill_color=OLIVE_BG, border_color=OLIVE_DARK)
textbox(s5, 9.15, 1.4, 3.7, 0.4, "Benefits:", size=18, bold=True, color=DARK_GRAY)
benefits = [
    "Near-zero marginal cost: sensors in phones, drivers already on roads",
    "Daily refresh vs annual surveys -- catches monsoon potholes immediately",
    "Composite severity score (IMU+Vision+Traffic) aligned to Indian standards",
    "MTTR dashboard -- the report that gets municipal cheques signed",
    "Open dataset contribution (SETU-IND-1) -- genuine research value",
    "SDK model: sell to 10 platforms, not 100M citizens",
]
txBox = s5.shapes.add_textbox(Inches(9.15), Inches(1.85), Inches(3.7), Inches(4.3))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(benefits):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    p.space_after = Pt(5)
    run = p.add_run()
    run.text = "  " + item
    run.font.size = Pt(11)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"


# ============ SLIDE 6: RESEARCH AND REFERENCES ============
print("Slide 6 - Research & References...")
s6 = prs.slides[5]

for sh in s6.shapes:
    if 'Oval' in sh.name and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = r.text.replace('Your Team Name', 'Team SETU')
    if sh.name == 'TextBox 8':
        clear_tf(sh.text_frame)
    if sh.name == 'Footer Placeholder 6':
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = '@SIH Idea submission - SETU'

refs = [
    "Sattar, Li & Chapman (2018) -- Road Surface Monitoring Using Smartphone Sensors: A Review, Sensors 18:3845",
    "Jan et al. (2023) -- Crowdsensing for Road Pavement Condition Monitoring, IEEE Access 11:133143",
    "Khandakar et al. (2025) -- Harnessing Smartphone Sensors for Enhanced Road Safety, Scientific Data 12:418",
    "MDPI Sensors (2020) 20(19):5564 -- Automated ML Approach for Road Pothole Detection Using Smartphone Sensors",
    "Arya, Maeda, Ghosh et al. (2022) -- RDD2022: Multi-national image dataset for Road Damage Detection, arXiv 2209.08538",
    "YOLOv8-PD, Nature Sci. Rep. (2024) -- 2.3M params, +1.4pp mAP while 74% baseline size",
    "MDPI Sensors 20(2):409 -- GPS localisation error 27-32m at highway speeds (critical design constraint)",
    "Celaya-Padilla et al. (2018) -- Speed Bump Detection Using Accelerometric Features, Sensors 18(2):443",
    "Nature Sci. Rep. (2026) -- Road roughness via smartphone + SVM (validates 50m segment analysis)",
    "MIT Carbin App -- 250,000+ miles collected; proves crowdsourced road data credible to agencies",
    "World Bank -- Road crashes cost India 3-5% GDP/year; ~150,000 deaths/year",
    "MoRTH (Lok Sabha tabled data) -- 9,438 pothole deaths 2020-2024, +53% rise in crashes",
    "NHAI Network Survey Vehicles -- deployed across 23 states, ~20,933 km of highways",
    "IRC:82-2015 -- Code of Practice for Maintenance of Bituminous Road Surfaces",
    "Digital Personal Data Protection Act, 2023 (DPDP Rules 13 Nov 2025) -- compliance framework",
]

txBox = s6.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.0), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, ref in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = "  " + ref
    run.font.size = Pt(11)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"


# ============ DELETE SLIDE 7 (Instructions) ============
print("Removing Slide 7 (Instructions)...")
last_sldId = prs.slides._sldIdLst[-1]
rId = last_sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
if rId is None:
    # Try without namespace
    for attr in last_sldId.attrib:
        if 'id' in attr.lower() and attr != 'id':
            rId = last_sldId.get(attr)
            break
if rId:
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(last_sldId)
    print("  Slide 7 removed.")
else:
    print("  Could not find relationship ID for slide 7, skipping removal.")


# ============ SAVE ============
print(f"\nSaving to: {OUT_PPTX}")
prs.save(OUT_PPTX)
print("PPTX saved successfully!")
print("Done!")
