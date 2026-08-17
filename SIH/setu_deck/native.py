"""
Native PowerPoint renderer for the SETU deck panels.

``PptxPanel`` implements the same drawing surface as ``layout.Panel`` but
emits real PowerPoint shapes and real text runs instead of rasterising to
a PNG. Because every composite in ``layout.py`` is expressed with the
primitives ``rrect / rect / circle / poly / arrow_* / raw_text``, the
panels in ``panels.py`` are reused byte-for-byte - only the back end
changes, so the exported PDF carries selectable, searchable text.

Geometry is identical to the bitmap version: the panel keeps its own
inch coordinate system with the origin at the top-left, and the panel is
placed on the slide at (ox, oy) inches.

Text metrics still come from matplotlib (same Segoe UI TTF that
PowerPoint uses), so wrapping and auto-shrinking are unchanged.
"""

import matplotlib.pyplot as plt

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

import layout
from layout import Panel, INK, INK_SOFT, WHITE, FONT_NAME


# ---------------------------------------------------------------- helpers
def rgb(c):
    """'#RRGGBB' (or an RGBColor) -> RGBColor."""
    if c is None:
        return None
    if isinstance(c, RGBColor):
        return c
    s = str(c).strip()
    if s.lower() in ("none", "transparent", ""):
        return None
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    return RGBColor.from_string(s.upper())


def is_none(c):
    return c is None or str(c).strip().lower() in ("none", "transparent", "")


ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT}

# vertical-centre of the glyph box, expressed in ems below the anchor y,
# so a matplotlib `va` maps onto a middle-anchored PowerPoint text frame
VA_SHIFT = {"top": 0.60, "center": 0.0, "center_baseline": 0.0,
            "baseline": -0.35, "bottom": -0.60}


class _AxProxy:
    """Stands in for `panel.ax` so `panel.ax.text(...)` calls in panels.py
    land on the native renderer."""

    def __init__(self, panel):
        self.p = panel

    def text(self, x, y, s, size=9, weight="normal", color=INK,
             ha="left", va="top", style="normal", zorder=6, **kw):
        self.p.raw_text(x, y, s, size=size, weight=weight, color=color,
                        ha=ha, va=va, style=style, z=zorder)

    def add_patch(self, *a, **k):        # pragma: no cover - guard rail
        raise NotImplementedError(
            "PptxPanel draws native shapes; use the panel primitives "
            "(rrect / rect / circle / poly / arrow_*) instead of add_patch.")


class PptxPanel(Panel):
    """Panel that draws straight onto a python-pptx slide."""

    def __init__(self, slide, ox, oy, w, h, dpi=200, bg=WHITE):
        super().__init__(w, h, dpi, bg)     # matplotlib fig = metrics only
        self.slide = slide
        self.ox, self.oy = ox, oy
        self.ax = _AxProxy(self)
        self._ops = []                     # (zorder, seq, callable)
        self._seq = 0
        if not is_none(bg):
            self.rect(0, 0, w, h, bg, z=-100)

    # ------------------------------------------------------- plumbing
    def _push(self, z, fn):
        self._ops.append((z, self._seq, fn))
        self._seq += 1

    def _X(self, x):
        return Inches(self.ox + x)

    def _Y(self, y):
        return Inches(self.oy + y)

    def _style(self, shp, fc, ec, lw, alpha=1.0):
        shp.shadow.inherit = False
        if is_none(fc):
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = rgb(fc)
            if alpha < 1.0:
                self._alpha(shp.fill.fore_color, alpha)
        if is_none(ec) or not lw:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = rgb(ec)
            shp.line.width = Pt(lw)
        return shp

    @staticmethod
    def _alpha(color, alpha):
        """Add <a:alpha/> to a solid fill colour."""
        srgb = color._xFill.find(qn("a:srgbClr"))
        if srgb is None:
            return
        el = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
        srgb.append(el)

    def _autoshape(self, kind, x, y, w, h):
        return self.slide.shapes.add_shape(kind, self._X(x), self._Y(y),
                                           Inches(w), Inches(h))

    def _freeform(self, pts, fc, ec=None, lw=1.0, alpha=1.0, close=True):
        b = self.slide.shapes.build_freeform(self._X(pts[0][0]),
                                             self._Y(pts[0][1]))
        b.add_line_segments([(self._X(px), self._Y(py)) for px, py in pts[1:]],
                            close=close)
        return self._style(b.convert_to_shape(), fc, ec, lw, alpha)

    # ------------------------------------------------------- primitives
    def rrect(self, x, y, w, h, fc=WHITE, ec=None, lw=1.1, r=0.07,
              z=2, alpha=1.0, shadow=False):
        if shadow:
            self.rrect(x + 0.035, y + 0.045, w, h, fc="#C9D4E2", ec=None,
                       r=r, z=z - 1, alpha=0.55)

        def draw():
            adj = 0.5 if min(w, h) <= 0 else min(max(r / min(w, h), 0.0), 0.5)
            shp = self._autoshape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            self._style(shp, fc, ec, lw, alpha)
            shp.adjustments[0] = adj
            shp.text_frame.word_wrap = False
        self._push(z, draw)

    def rect(self, x, y, w, h, fc, alpha=1.0, z=2):
        def draw():
            shp = self._autoshape(MSO_SHAPE.RECTANGLE, x, y, w, h)
            self._style(shp, fc, None, 0, alpha)
        self._push(z, draw)

    def circle(self, cx, cy, r, fc, ec=None, lw=1.2, z=3):
        def draw():
            shp = self._autoshape(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r)
            self._style(shp, fc, ec, lw)
        self._push(z, draw)

    def poly(self, pts, fc, ec=None, lw=1.0, z=2, alpha=1.0):
        self._push(z, lambda: self._freeform(list(pts), fc, ec, lw, alpha))

    def arrow_down(self, cx, y0, y1, color=layout.MUTED, w=0.055, head=0.13):
        neck = y1 - head
        self._push(4, lambda: self._freeform(
            [(cx - w / 2, y0), (cx + w / 2, y0), (cx + w / 2, neck),
             (cx + w * 1.8, neck), (cx, y1), (cx - w * 1.8, neck),
             (cx - w / 2, neck)], color, None, 0))

    def arrow_right(self, x0, cy, x1, color=layout.MUTED, w=0.05, head=0.12):
        neck = x1 - head
        self._push(4, lambda: self._freeform(
            [(x0, cy - w / 2), (neck, cy - w / 2), (neck, cy - w * 1.8),
             (x1, cy), (neck, cy + w * 1.8), (neck, cy + w / 2),
             (x0, cy + w / 2)], color, None, 0))

    def curved_arrow(self, x0, y0, x1, y1, rad=0.25, color=layout.MUTED,
                     lw=1.5, dashed=True, z=4, scale=11):
        """Quadratic-Bezier stand-in for matplotlib's arc3 connector."""
        mx, my = (x0 + x1) / 2, (y0 + y1) / 2
        dx, dy = x1 - x0, y1 - y0
        ctrl = (mx - dy * rad, my + dx * rad)          # perpendicular offset
        pts = []
        for i in range(25):
            t = i / 24.0
            u = 1 - t
            pts.append((u * u * x0 + 2 * u * t * ctrl[0] + t * t * x1,
                        u * u * y0 + 2 * u * t * ctrl[1] + t * t * y1))

        def draw():
            b = self.slide.shapes.build_freeform(self._X(pts[0][0]),
                                                 self._Y(pts[0][1]))
            b.add_line_segments([(self._X(px), self._Y(py))
                                 for px, py in pts[1:]], close=False)
            shp = b.convert_to_shape()
            shp.shadow.inherit = False
            shp.fill.background()
            shp.line.color.rgb = rgb(color)
            shp.line.width = Pt(lw)
            if dashed:
                ln = shp.line._get_or_add_ln()
                d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
                ln.append(d)
            # arrow head
            hx, hy = pts[-1]
            px, py = pts[-3]
            vx, vy = hx - px, hy - py
            n = (vx * vx + vy * vy) ** 0.5 or 1.0
            vx, vy = vx / n, vy / n
            hl = 0.10
            bx, by = hx - vx * hl, hy - vy * hl
            self._freeform([(hx, hy), (bx - vy * hl * 0.45,
                                       by + vx * hl * 0.45),
                            (bx + vy * hl * 0.45, by - vx * hl * 0.45)],
                           color, None, 0)
        self._push(z, draw)

    # ------------------------------------------------------- text
    def raw_text(self, x, y, s, size=9, weight="normal", color=INK_SOFT,
                 ha="left", va="top", style="normal", z=6):
        if s is None or str(s) == "":
            return
        s = str(s)
        bold = weight in ("bold", "semibold", "heavy", "black") or \
            (isinstance(weight, (int, float)) and weight >= 600)
        italic = style == "italic"
        em = size / 72.0
        cy = y + VA_SHIFT.get(va, 0.60) * em
        tw = self.text_w(s, size, "bold" if bold else "normal", style)
        boxw = tw + 0.12
        boxh = max(em * 1.9, 0.14)
        if ha == "center":
            left = x - boxw / 2
        elif ha == "right":
            left = x - boxw + 0.06
        else:
            left = x - 0.055                  # frame inset vs ink start

        def draw():
            tb = self.slide.shapes.add_textbox(
                self._X(left), self._Y(cy - boxh / 2), Inches(boxw),
                Inches(boxh))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = ALIGN.get(ha, PP_ALIGN.LEFT)
            r = p.add_run()
            r.text = s
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.name = FONT_NAME
            f.color.rgb = rgb(color)
        self._push(z, draw)

    # ------------------------------------------------------- output
    def flush(self):
        """Emit every buffered op, lowest zorder first (stable)."""
        for _, _, fn in sorted(self._ops, key=lambda o: (o[0], o[1])):
            fn()
        n = len(self._ops)
        self._ops = []
        return n

    def save(self, path=None):
        """Same call signature as Panel.save so panels.py is untouched."""
        n = self.flush()
        plt.close(self.fig)
        self.shape_count = n
        return path


# ---------------------------------------------------------------- driver
def render(slide, panel_fn, x, y):
    """Run a panels.py builder with the native renderer swapped in.

    The panel function keeps full control of its own size; only the
    placement on the slide comes from (x, y) in inches.
    """
    import panels

    made = {}

    def factory(w, h, dpi=200, bg=WHITE):
        p = PptxPanel(slide, x, y, w, h, dpi, bg)
        made["p"] = p
        return p

    prev = panels.Panel
    panels.Panel = factory
    try:
        panel_fn()
    finally:
        panels.Panel = prev
    p = made.get("p")
    return getattr(p, "shape_count", 0)
