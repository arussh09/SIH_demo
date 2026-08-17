"""
Build the SETU SIH 2026 Idea Presentation.

Uses the OFFICIAL SIH template PPTX untouched (branding, title font,
footer bar, slide numbers, team oval) and drops in rebuilt diagram
panels + a native, click-through reference slide.
"""

import os
import shutil
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

import panels
import native
import make_logo

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, "Sih-ppt-template-2025-pdf-download.pptx")
OUT_PPTX = os.path.join(ROOT, "SETU_SIH2026_Idea_Presentation.pptx")
OUT_PDF = os.path.join(ROOT, "SETU_SIH2026_Idea_Presentation.pdf")
ALIAS_PDF = os.path.join(ROOT, "SETU.pdf")
TARGET_PDF = os.path.join(ROOT, "SETU_SIH2026.pdf")
LOGO_BYTES = 12933          # size of the template's 2025 logo bitmap

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A}

# palette (mirrors panels/layout)
INK = RGBColor(0x0B, 0x1B, 0x33)
INK_SOFT = RGBColor(0x3C, 0x4C, 0x63)
MUTED = RGBColor(0x6B, 0x7A, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINK = RGBColor(0x12, 0x5E, 0xC4)
NAVY = RGBColor(0x12, 0x29, 0x4B)

CLR = {
    "blue":   (RGBColor(0x15, 0x65, 0xD8), RGBColor(0xE8, 0xF1, 0xFE)),
    "teal":   (RGBColor(0x0E, 0x95, 0x94), RGBColor(0xE3, 0xF6, 0xF5)),
    "green":  (RGBColor(0x17, 0x8A, 0x4C), RGBColor(0xE6, 0xF6, 0xEC)),
    "amber":  (RGBColor(0xD9, 0x8A, 0x00), RGBColor(0xFF, 0xF5, 0xE1)),
    "red":    (RGBColor(0xD9, 0x30, 0x25), RGBColor(0xFD, 0xEC, 0xEA)),
    "purple": (RGBColor(0x6C, 0x3B, 0xD1), RGBColor(0xF0, 0xEA, 0xFE)),
}


# ---------------------------------------------------------------- helpers
def shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def clear_tf(tf):
    body = tf._txBody
    paras = body.findall("a:p", NS)
    for p in paras[1:]:
        body.remove(p)
    for child in list(paras[0]):
        if etree.QName(child.tag).localname != "pPr":
            paras[0].remove(child)


def kill(shape):
    """Remove a shape from its slide."""
    shape._element.getparent().remove(shape._element)


def add_run(p, text, size, bold=False, color=INK_SOFT, italic=False,
            font="Calibri", url=None, underline=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    f.underline = underline
    if url:
        r.hyperlink.address = url
    return r


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def rounded(slide, x, y, w, h, fill, line=None, lw=1.0, adj=0.14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                 Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.adjustments[0] = adj
    shp.text_frame.word_wrap = True
    return shp


def set_title(slide, x=1.85, y=0.05, w=8.70, h=1.05, size=27, text=None):
    """Keep the template's title font/colour, just fit it between the
    team oval and the SIH logo."""
    t = shape_by_name(slide, "Title 1")
    if t is None:
        return
    t.left, t.top, t.width, t.height = (Inches(x), Inches(y), Inches(w),
                                        Inches(h))
    tf = t.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.text = r.text.replace("\x0b", " ").strip()
            if text is not None:
                r.text = text if first else ""
                first = False
            r.font.size = Pt(size)

def brand_slide(slide, team="localhost"):
    """Team oval + footer, then drop the template's instruction text box."""
    for s in slide.shapes:
        if s.name.startswith("Oval") and s.has_text_frame:
            tf = s.text_frame
            clear_tf(tf)
            tf.word_wrap = False            # keep "localhost" on one line
            tf.margin_left = tf.margin_right = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            # styled like a terminal prompt - the team is called localhost
            add_run(p, team, 12, bold=True, color=NAVY, font="Consolas")
            add_run(p, "_", 12, bold=True, color=LINK, font="Consolas")
        if s.name == "Footer Placeholder 6":
            for p in s.text_frame.paragraphs:
                for i, r in enumerate(p.runs):
                    r.text = "@SIH 2026 Idea Submission - Team localhost" if i == 0 else ""
    box = shape_by_name(slide, "TextBox 8")
    if box is not None:
        kill(box)


def put_panel(slide, panel_fn, x=0.30, y=1.20):
    """Draw a panels.py diagram as native PowerPoint shapes + real text
    (no bitmap), anchored at (x, y) inches."""
    return native.render(slide, panel_fn, x, y)


def swap_logo(slide, logo):
    """Drop the template's 2025 logo bitmap and put the 2026 one in its
    place (identical rectangle, identical aspect ratio)."""
    for s in list(slide.shapes):
        if s.shape_type == 13 and len(s.image.blob) == LOGO_BYTES:
            box = (s.left, s.top, s.width, s.height)
            kill(s)
            slide.shapes.add_picture(logo, box[0], box[1], box[2], box[3])
            return True
    return False


# ---------------------------------------------------------------- content
METADATA = [
    ("Problem Statement ID", "SIH 2026 (to be filled)"),
    ("Problem Statement Title", "Crowdsourced Road Defect Intelligence for India"),
    ("Theme", "Smart Automation"),
    ("PS Category", "Software"),
    ("Team ID", "To be filled"),
    ("Team Name", "localhost"),
]

RESEARCH = [
    ("Sattar, Li & Chapman - Sensors 18:3845 (2018)",
     "A phone's motion sensor can spot road damage; sampling rate is the "
     "real limit, so we sample at 100 Hz.",
     "mdpi.com/1424-8220/18/11/3845",
     "https://www.mdpi.com/1424-8220/18/11/3845", "blue"),
    ("Wu, Wang, Hu et al. - Sensors 20:5564 (2020)",
     "88.5% precision and 75% recall per 10 m window - the base numbers our "
     "four-phone agreement math is built on.",
     "mdpi.com/1424-8220/20/19/5564",
     "https://www.mdpi.com/1424-8220/20/19/5564", "teal"),
    ("Sensors 20:409 (2020) - participatory sensing",
     "Measured GPS error of 27-32 m at speed, which is why we group "
     "inside 20 m after map-matching.",
     "mdpi.com/1424-8220/20/2/409",
     "https://www.mdpi.com/1424-8220/20/2/409", "amber"),
    ("Arya, Maeda, Ghosh et al. - RDD2022 (arXiv)",
     "47,420 labelled road images from 6 countries including India - our "
     "training base for the vision model.",
     "arxiv.org/abs/2209.08538",
     "https://arxiv.org/abs/2209.08538", "red"),
]

STANDARDS = [
    ("MoRTH - Road Accidents in India",
     "9,438 pothole-related deaths in 2020-24, up 53% in five years.",
     "morth.nic.in/road-accident-in-india",
     "https://morth.nic.in/road-accident-in-india", "red"),
    ("NHAI Network Survey Vehicles",
     "The survey fleet covers ~20,933 km of highway about once a year, at "
     "high cost. We cover every lane driven.",
     "nhai.gov.in",
     "https://nhai.gov.in", "teal"),
    ("TrafficSense - Mohan, Padmanabhan & Ramjee (MSR, 2008)",
     "Phones detected potholes and bumps on Bangalore roads 17 years ago - "
     "the idea is field-proven in Indian traffic.",
     "microsoft.com/research - TrafficSense (MSR-TR-2008-59)",
     "https://www.microsoft.com/en-us/research/publication/trafficsense-rich-monitoring-of-road-and-traffic-conditions-using-mobile-smartphones/",
     "purple"),
    ("NITI Aayog Frontier Tech Hub - RoadMetrics",
     "An Indian phone-camera road-mapping startup: 50,000+ km mapped, "
     "adopted by Chennai. The market is real; we make it free.",
     "frontiertech.niti.gov.in - RoadMetrics story",
     "https://frontiertech.niti.gov.in/story/ai-driven-road-management-enhancing-indias-infrastructure-with-roadmetrics/",
     "blue"),
]


# ---------------------------------------------------------------- slides
def build():
    logo = make_logo.build()

    print("Loading official SIH template...")
    prs = Presentation(TEMPLATE)
    s1, s2, s3, s4, s5, s6 = [prs.slides[i] for i in range(6)]

    # ---------------- slide 1 : title ----------------
    t = shape_by_name(s1, "Title 7")
    for p in t.text_frame.paragraphs:
        for r in p.runs:
            r.text = r.text.replace("2025", "2026")

    # the SETU expansion and the tagline both move off the title slide -
    # the full form now heads slide 2, and the KPI numbers move into the
    # problem block there.
    sub = shape_by_name(s1, "Subtitle 3")
    if sub is not None:
        kill(sub)

    meta = shape_by_name(s1, "TextBox 9")
    meta.left, meta.top = Inches(0.55), Inches(1.62)
    meta.width, meta.height = Inches(6.55), Inches(5.30)
    tf = meta.text_frame
    clear_tf(tf)
    for i, (k, v) in enumerate(METADATA):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        p.line_spacing = 1.06
        pPr = p._p.get_or_add_pPr()          # hanging indent for wraps
        pPr.set("marL", str(int(Inches(0.42))))
        pPr.set("indent", str(int(-Inches(0.42))))
        add_run(p, "\u2022  " + k + "- ", 23, bold=False, color=INK)
        add_run(p, v, 23, bold=True, color=INK)

    print("Drawing panels as native shapes + text...")

    # ---------------- slides 2-5 : panels ----------------
    titles = {
        id(s2): "SETU - Sensor-Enabled Tracking of Urban-road-damage",
    }
    sizes = {id(s2): 25}
    for i, (slide, fn) in enumerate(((s2, panels.slide2), (s3, panels.slide3),
                                     (s4, panels.slide4),
                                     (s5, panels.slide5)), start=2):
        set_title(slide, text=titles.get(id(slide)),
                  size=sizes.get(id(slide), 34))
        brand_slide(slide)
        n = put_panel(slide, fn)
        print(f"  slide {i} panel     : {n} native shapes")

    for slide in (s1, s2, s3, s4, s5, s6):
        swap_logo(slide, logo)

    # ---------------- slide 6 : references (native, clickable) --------
    # Deliberately plain: no coloured cards, no tints - just typed
    # bullets the way a person would list sources in PowerPoint.
    set_title(s6, size=34)
    brand_slide(s6)
    cols = [(0.42, "Peer-reviewed foundation", RESEARCH),
            (6.72, "Government data and Indian field proof", STANDARDS)]
    cw = 6.05
    for x0, header, items in cols:
        tfh = textbox(s6, x0, 1.14, cw, 0.40)
        ph = tfh.paragraphs[0]
        add_run(ph, header, 17, bold=True, color=INK)

        y = 1.66
        for title, why, label, url, key in items:
            tf = textbox(s6, x0, y, cw, 1.10)
            p = tf.paragraphs[0]
            p.space_after = Pt(3)
            add_run(p, "\u2022  " + title, 14, bold=True, color=INK)
            p = tf.add_paragraph()
            p.space_after = Pt(3)
            add_run(p, "     " + why, 12.5, color=INK_SOFT)
            p = tf.add_paragraph()
            add_run(p, "     " + label, 12, color=LINK, url=url,
                    underline=True)
            y += 1.26

    # ---------------- drop the instruction slide ----------------
    sldIdLst = prs.slides._sldIdLst
    last = list(sldIdLst)[-1]
    rId = last.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    sldIdLst.remove(last)

    prs.save(OUT_PPTX)
    print("PPTX ->", OUT_PPTX)
    return OUT_PPTX


def to_pdf():
    ps = (
        "$ppt = New-Object -ComObject PowerPoint.Application;"
        f"$d = $ppt.Presentations.Open('{OUT_PPTX}', $true, $false, $false);"
        f"$d.SaveAs('{OUT_PDF}', 32);"
        "$d.Close(); $ppt.Quit();"
    )
    subprocess.run(["powershell", "-NoProfile", "-Command", ps], check=True)
    add_pdf_links()
    shutil.copyfile(OUT_PDF, ALIAS_PDF)
    shutil.copyfile(OUT_PDF, TARGET_PDF)
    print("PDF  ->", OUT_PDF)
    print("PDF  ->", ALIAS_PDF)
    print("PDF  ->", TARGET_PDF)


def add_pdf_links():
    """PowerPoint's PDF writer drops run-level hyperlinks, so re-attach
    them onto the reference slide by locating the printed URL labels."""
    import fitz
    doc = fitz.open(OUT_PDF)
    page = doc[-1]                                  # references slide
    added = 0
    for _, _, label, url, _ in RESEARCH + STANDARDS:
        hits = page.search_for(label)
        if not hits:
            hits = page.search_for(label.split(" - ")[0])
        if not hits:
            print("  !! link label not found:", label)
            continue
        rect = hits[0]
        for extra in hits[1:]:
            rect |= extra
        page.insert_link({"kind": fitz.LINK_URI, "from": rect, "uri": url})
        added += 1
    doc.saveIncr()
    doc.close()
    print(f"Re-attached {added} clickable reference links in the PDF")


def preview():
    import fitz
    doc = fitz.open(OUT_PDF)
    outdir = os.path.join(ROOT, "preview")
    os.makedirs(outdir, exist_ok=True)
    print("Text coverage in the exported PDF:")
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=110)
        pix.save(os.path.join(outdir, f"slide_{i + 1}.png"))
        txt = page.get_text().strip()
        print(f"  page {i + 1}: {len(txt):>5} selectable chars, "
              f"{len(page.get_text('words')):>4} words, "
              f"{len(page.get_images(full=True))} raster image(s)")
    print(f"Rendered {len(doc)} preview page(s) -> {outdir}")
    doc.close()


if __name__ == "__main__":
    build()
    if "--pdf" in sys.argv or len(sys.argv) == 1:
        to_pdf()
        preview()
