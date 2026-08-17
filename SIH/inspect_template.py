"""Inspect the template PPTX to understand its structure."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r"d:\Testing\Research\SIH\Sih-ppt-template-2025-pdf-download.pptx")

print(f"Slide width: {prs.slide_width} EMU = {prs.slide_width / 914400:.2f} inches")
print(f"Slide height: {prs.slide_height} EMU = {prs.slide_height / 914400:.2f} inches")
print(f"Number of slides: {len(prs.slides)}")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_idx + 1} - Layout: '{slide.slide_layout.name}'")
    print(f"  Shapes: {len(slide.shapes)}")
    
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"\n  [{shape_idx}] '{shape.name}' type={shape.shape_type}")
        print(f"      pos=({shape.left/914400:.2f}, {shape.top/914400:.2f}) in")
        print(f"      size=({shape.width/914400:.2f}, {shape.height/914400:.2f}) in")
        
        if hasattr(shape, "text") and shape.text:
            text_preview = shape.text[:150].replace('\n', ' | ')
            print(f"      TEXT: '{text_preview}'")
        
        if shape.shape_type == 13:  # Picture
            print(f"      IMAGE: {shape.image.content_type}")
        
        # Check for group shapes
        if shape.shape_type == 6:  # Group
            print(f"      GROUP with {len(shape.shapes)} sub-shapes")

print("\n\nDone!")
